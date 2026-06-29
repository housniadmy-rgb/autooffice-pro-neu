"""Build Investor_Pitch_Deck.pptx — ~12 professional slides for buyer outreach."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUT = Path(__file__).parent / "Investor_Pitch_Deck.pptx"

# Brand colors
PRIMARY = RGBColor(0x1F, 0x49, 0x7D)
ACCENT = RGBColor(0x2A, 0x80, 0xD9)
LIGHT_BG = RGBColor(0xF6, 0xF8, 0xFB)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill_solid(shape, color)
    return shape

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_GREY, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def add_footer(slide, page_num, total=12):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(7), Inches(0.3),
             "PraxisOnline24 — Investor Pitch Deck",
             size=9, color=MID_GREY)
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{page_num} / {total}", size=9, color=MID_GREY, align=PP_ALIGN.RIGHT)

# ────────────────────────────────────────────────────────────────────
# Slide 1: Title
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
add_text(s, Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.2),
         "PraxisOnline24", size=66, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(0.6),
         "Production-ready, 12-language SaaS for medical practices",
         size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.5),
         "Ihre Praxis. Online. Rund um die Uhr.",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
         "Pre-Revenue Asset for Acquisition  |  www.praxisonline24.com",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ────────────────────────────────────────────────────────────────────
# Slide 2: Problem
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "The Problem", size=36, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
         "Small practices have no good way to accept bookings 24/7 — and no way to do it cheaply.",
         size=18, italic=True, color=MID_GREY)
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(4),
            ["Phone-only booking leaves money on the table — patients call competitors when lines are busy.",
             "Cancellations are not automatically refilled, wasting practitioner time.",
             "Generic tools (Calendly + spreadsheets) cannot handle multi-practitioner availability, invoicing, waitlists.",
             "Marketplace incumbents (Doctolib, Jameda) are expensive and force practices into the marketplace they do not control.",
             "Non-English markets are deeply under-served — almost no competitor ships in 12 languages."],
            size=18)
add_footer(s, 2)

# ────────────────────────────────────────────────────────────────────
# Slide 3: Solution
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "The Solution", size=36, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
         "An integrated, white-label-friendly platform — 12 languages, self-hostable, privacy-by-design.",
         size=18, italic=True, color=MID_GREY)
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(4.4),
            ["3-step public booking flow with ICS export and one-token cancellation",
             "Multi-practitioner calendar with availability windows + conflict detection",
             "PDF invoicing with VAT + line items",
             "Waitlist with cron-driven auto-offer engine (4h tokens)",
             "Review collection + moderation",
             "AI Practice Manager: 14-section daily dashboard, traffic-light recommendations",
             "12 languages (DE/EN/FR/ES/PT/AR-RTL/RU/ZH/HI/TH/TR/ID)",
             "30-day trial → automatic pause cycle, no manual ops"],
            size=16)
add_footer(s, 3)

# ────────────────────────────────────────────────────────────────────
# Slide 4: Target Audience
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Target Audience", size=36, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
         "Owner-operated practice, 3–10 practitioners, >50 appointments/week.",
         size=18, italic=True, color=MID_GREY)
add_bullets(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(4.4),
            ["Doctors, therapists, alternative practitioners",
             "Physiotherapy, chiropractic, osteopathy",
             "Veterinary practices",
             "Beauty, wellness, cosmetics",
             "Driving schools",
             "Tax and legal consultancies",
             "Independent coaches and therapists",
             "Budget profile: €300–€1,500/year for booking software"],
            size=17)
add_footer(s, 4)

# ────────────────────────────────────────────────────────────────────
# Slide 5: Market
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Market Size", size=36, bold=True, color=PRIMARY)

# Key stats boxes
def stat_box(x, top, label, value, subtitle):
    add_rect(s, x, top, Inches(3.8), Inches(2.0), LIGHT_BG)
    add_text(s, x + Inches(0.1), top + Inches(0.1), Inches(3.6), Inches(0.4),
             label, size=12, color=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), top + Inches(0.5), Inches(3.6), Inches(0.9),
             value, size=36, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), top + Inches(1.5), Inches(3.6), Inches(0.4),
             subtitle, size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

stat_box(Inches(0.6),  Inches(1.8), "GERMANY DOCTOR PRACTICES", "~204,000", "KBV public registry")
stat_box(Inches(4.7),  Inches(1.8), "EU MEDICAL + PARAMEDICAL", ">1,000,000", "Eurostat / national chambers")
stat_box(Inches(8.8),  Inches(1.8), "LANGUAGES SHIPPED", "12", "incl. RTL Arabic")

add_text(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5),
         "Competitive landscape", size=22, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2),
            ["Doctolib — DACH leader, marketplace play, expensive for the practice",
             "samedi — feature-rich, enterprise-leaning pricing",
             "Jameda — review-led; booking is secondary",
             "Calendly + manual workflows — low cost, missing invoicing / waitlist / multi-practitioner"],
            size=14)
add_footer(s, 5)

# ────────────────────────────────────────────────────────────────────
# Slide 6: Product Demo Highlights
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Product — What's Live Today", size=36, bold=True, color=PRIMARY)

# Two-column features
add_text(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
         "Patient surface", size=20, bold=True, color=ACCENT)
add_bullets(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4),
            ["3-step booking flow",
             "ICS calendar download",
             "Cancel via one-time token",
             "Patient portal lookup",
             "Public review form",
             "Waitlist signup + auto-offer"],
            size=15)

add_text(s, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5),
         "Admin surface", size=20, bold=True, color=ACCENT)
add_bullets(s, Inches(6.8), Inches(2.1), Inches(5.5), Inches(4.5),
            ["Calendar (day/week/month) + archive filter",
             "Practitioner CRUD + weekly availability",
             "Invoices with VAT + PDF download",
             "Review moderation",
             "Recipe print log (audit-only)",
             "Activity log (every admin action)",
             "AI Practice Manager (14-section dashboard)"],
            size=15)
add_footer(s, 6)

# ────────────────────────────────────────────────────────────────────
# Slide 7: Tech Stack
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Tech Stack", size=36, bold=True, color=PRIMARY)

stack = [
    ("Runtime", "Node.js ≥ 18 LTS"),
    ("Framework", "Express 4.18 (monolith, 17 modular routes)"),
    ("Database", "SQLite via sql.js — easy backup, single file"),
    ("Auth", "express-session + bcrypt cost 12 + role-based middleware"),
    ("Frontend", "Vanilla HTML5 / CSS3 / ES modules — no framework lock-in"),
    ("PDF", "pdfkit"),
    ("Scheduling", "node-cron (8 jobs)"),
    ("E-mail", "nodemailer (SMTP) + Brevo API"),
    ("i18n", "12 JSON catalogs, RTL Arabic supported"),
    ("Hosting", "Render.com Frankfurt (render.yaml committed)"),
    ("Front", "Cloudflare (TLS + CDN)"),
    ("Backups", "Daily SQLite snapshot, 7-day rolling"),
]
# Two columns
col_x = [Inches(0.8), Inches(6.9)]
col_y_start = Inches(1.7)
row_h = Inches(0.42)
for i, (k, v) in enumerate(stack):
    col = i % 2
    row = i // 2
    y = col_y_start + row_h * row
    add_text(s, col_x[col], y, Inches(2.1), row_h, k, size=14, bold=True, color=ACCENT)
    add_text(s, col_x[col] + Inches(2.1), y, Inches(3.5), row_h, v, size=13, color=DARK_GREY)
add_footer(s, 7)

# ────────────────────────────────────────────────────────────────────
# Slide 8: Competitive Advantages
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "What Makes This Different", size=36, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5),
            ["12 languages, including RTL Arabic — rare in this price band",
             "Self-hostable on a $7/month Render plan",
             "Privacy-by-design data model (no health data ever stored)",
             "30-day auto-anonymisation cron, daily SQLite backup, full audit log",
             "Multi-tenant from day one — practice_id everywhere",
             "AI Practice Manager wired into the admin home (rule engine; LLM-pluggable)",
             "No framework drift risk — vanilla HTML/JS frontend",
             "Brand + domain + Cloudflare zone + render.yaml all included"],
            size=17)
add_footer(s, 8)

# ────────────────────────────────────────────────────────────────────
# Slide 9: Project Statistics
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Project Statistics", size=36, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
         "All numbers counted from the current repo state.",
         size=14, italic=True, color=MID_GREY)

def stat_card(x, y, top_text, big_text, sub_text):
    add_rect(s, x, y, Inches(3.0), Inches(1.6), LIGHT_BG)
    add_text(s, x + Inches(0.05), y + Inches(0.1), Inches(2.9), Inches(0.4),
             top_text, size=10, color=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), y + Inches(0.45), Inches(2.9), Inches(0.7),
             big_text, size=26, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), y + Inches(1.18), Inches(2.9), Inches(0.4),
             sub_text, size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

cards = [
    ("CODEBASE",      "~17.2k", "JS server LoC"),
    ("FRONTEND",      "7.4k",   "HTML LoC (33 pages)"),
    ("STYLING",       "902",    "CSS LoC"),
    ("I18N",          "12",     "languages, 906 catalog LoC"),
    ("ROUTES",        "17",     "modular Express files"),
    ("DATABASE",      "18",     "SQL tables"),
    ("AUTOMATION",    "8",      "node-cron jobs"),
    ("QA CHECKLIST",  "187",    "manual verification points"),
]
positions = [
    (Inches(0.7),  Inches(2.0)), (Inches(3.85), Inches(2.0)),
    (Inches(7.0),  Inches(2.0)), (Inches(10.15),Inches(2.0)),
    (Inches(0.7),  Inches(3.9)), (Inches(3.85), Inches(3.9)),
    (Inches(7.0),  Inches(3.9)), (Inches(10.15),Inches(3.9)),
]
for (top, big, sub), (x, y) in zip(cards, positions):
    stat_card(x, y, top, big, sub)
add_footer(s, 9)

# ────────────────────────────────────────────────────────────────────
# Slide 10: Roadmap (post-acquisition)
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Recommended Buyer Roadmap", size=36, bold=True, color=PRIMARY)

# Three-column timeline
columns = [
    ("Days 1–30",   "Stabilize", [
        "Wipe test data, rotate secrets",
        "Re-enable Stripe checkout",
        "Add CSP, CSRF, 2FA",
        "Sentry / error tracking",
        "Lawyer review of German legal pages"]),
    ("Days 31–90",  "Harden",    [
        "Migrate SQLite → PostgreSQL",
        "Replace 'moment' → 'date-fns'",
        "~30 Playwright E2E specs",
        "Native review of RU/ZH/HI/TH",
        "Launch in 1–2 verticals"]),
    ("Months 4–12", "Scale",     [
        "First 50–200 paying practices",
        "Vertical-specific UI variants",
        "Patient mobile app (optional)",
        "Real LLM in AI Practice Manager",
        "Public review aggregator → SEO"]),
]
col_w = Inches(4.0)
gap = Inches(0.2)
start_x = Inches(0.6)
top_y = Inches(1.6)
for i, (phase, title, items) in enumerate(columns):
    x = start_x + (col_w + gap) * i
    add_rect(s, x, top_y, col_w, Inches(0.6), PRIMARY)
    add_text(s, x, top_y + Inches(0.1), col_w, Inches(0.4), phase,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, top_y + Inches(0.7), col_w, Inches(0.5), title,
             size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.15), top_y + Inches(1.4), col_w - Inches(0.3),
                Inches(4.0), items, size=13)
add_footer(s, 10)

# ────────────────────────────────────────────────────────────────────
# Slide 11: Valuation
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, Inches(0.4), prs.slide_height, PRIMARY)
add_text(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
         "Valuation", size=36, bold=True, color=PRIMARY)
add_text(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
         "Methodology: replacement cost + comparable pre-revenue SaaS sales. "
         "Full reasoning in Pricing_Strategy.md.",
         size=14, italic=True, color=MID_GREY)

# Three valuation tiers as boxes
def val_box(x, top, label, range_text, list_text, accent_color):
    add_rect(s, x, top, Inches(3.8), Inches(3.0), LIGHT_BG)
    add_rect(s, x, top, Inches(3.8), Inches(0.5), accent_color)
    add_text(s, x, top + Inches(0.07), Inches(3.8), Inches(0.4),
             label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), top + Inches(0.8), Inches(3.4), Inches(0.8),
             range_text, size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), top + Inches(1.8), Inches(3.4), Inches(1.0),
             list_text, size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)

val_box(Inches(0.6),  Inches(2.1), "QUICK SALE",
        "$8k – $15k", "10–15% of replacement cost\nAs-is, 7-day close, no support",
        RGBColor(0xC0, 0x39, 0x2B))
val_box(Inches(4.7),  Inches(2.1), "FAIR MARKET VALUE",
        "$18k – $32k", "25–35% of replacement cost\nList $24,500\n4–8 wk close, 30-day support",
        RGBColor(0x1F, 0x49, 0x7D))
val_box(Inches(8.8),  Inches(2.1), "OPTIMISTIC",
        "$35k – $55k", "40–60% of replacement cost\nStrategic buyer (12-lang i18n)\n60-day paid consulting",
        RGBColor(0x27, 0xAE, 0x60))

add_text(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(0.6),
         "Replacement cost (independent estimate): €75k–€120k  /  ~$80k–$130k",
         size=15, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
add_footer(s, 11)

# ────────────────────────────────────────────────────────────────────
# Slide 12: Call To Action
# ────────────────────────────────────────────────────────────────────
s = add_blank_slide()
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
add_text(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.9),
         "Ready to make an offer?",
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(0.6),
         "All due-diligence materials are in the sale/ folder.",
         size=20, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.5),
         "Demo: https://www.praxisonline24.com  |  Listing: see Acquire_Listing.md",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.0),
         "Next steps:\n"
         "1.   Sign NDA  →  read-only repo invite\n"
         "2.   Walk the Demo_Guide.md  →  27 minutes hands-on\n"
         "3.   Submit LOI via Acquire.com escrow\n"
         "4.   ~10 days to full handover incl. domain EPP transfer",
         size=18, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

prs.save(OUT)
print(f"Written: {OUT}")
print(f"Slides: {len(prs.slides)}")
